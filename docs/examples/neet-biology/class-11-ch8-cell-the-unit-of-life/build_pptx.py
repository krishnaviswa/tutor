"""Build the Class 11 NCERT Chapter 8 classroom PowerPoint (16:9)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
IMG = ROOT / "images"
OUT = ROOT / "Cell-The-Unit-of-Life-Class11.pptx"

TEAL = RGBColor(0x0D, 0x73, 0x77)
TEAL_D = RGBColor(0x08, 0x4E, 0x52)
INK = RGBColor(0x1C, 0x24, 0x2A)
CREAM = RGBColor(0xFF, 0xFB, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x8C, 0x30)
GOLD = RGBColor(0xF5, 0xC4, 0x48)
SOFT = RGBColor(0xE8, 0xF8, 0xF5)
TRAP_BG = RGBColor(0xFF, 0xF3, 0xD6)
PINK = RGBColor(0xC4, 0x3C, 0x6A)

W = Inches(13.333)
H = Inches(7.5)


def pick(*names: str) -> str | None:
    for name in names:
        path = IMG / name
        if path.exists() and path.stat().st_size > 800:
            return str(path)
    return None


def rgb_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def line(shape, color: RGBColor, pt: float = 1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def _set_run(run, text, size, bold=False, color=INK, italic=False):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", "Calibri")


def box(slide, l, t, w, h, fill=None, outline=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.08
    if fill is None:
        sh.fill.background()
    else:
        rgb_fill(sh, fill)
    if outline:
        line(sh, outline, 1.25)
    else:
        sh.line.fill.background()
    return sh


def tb(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(l, t, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        sh.text_frame._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    _set_run(run, text, size, bold, color)
    return sh


def bullets(slide, l, t, w, h, items, size=16, color=INK, bold_first=False, spacing=6):
    sh = slide.shapes.add_textbox(l, t, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        p.level = 0
        p.clear()
        run = p.add_run()
        _set_run(run, item, size, bold=(bold_first and i == 0), color=color)
    return sh


def notes(slide, text: str):
    ns = slide.notes_slide
    ns.notes_text_frame.text = text


def header(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.78))
    rgb_fill(bar, TEAL)
    bar.line.fill.background()
    tb(slide, Inches(0.35), Inches(0.12), Inches(10.5), Inches(0.55), title, 26, True, WHITE)
    if subtitle:
        tb(slide, Inches(9.4), Inches(0.22), Inches(3.6), Inches(0.4), subtitle, 12, False, GOLD, PP_ALIGN.RIGHT)


def footer(slide, n: int, total: int = 21):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.18), W, Inches(0.32))
    rgb_fill(bar, TEAL_D)
    bar.line.fill.background()
    tb(slide, Inches(0.3), Inches(7.18), Inches(9), Inches(0.3), "Class 11 Biology  |  Chapter 8  |  NCERT", 11, False, WHITE, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    tb(slide, Inches(10.5), Inches(7.18), Inches(2.5), Inches(0.3), f"{n} / {total}", 11, False, WHITE, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)


def picture(slide, path, l, t, w, h):
    if not path:
        box(slide, l, t, w, h, SOFT, TEAL)
        tb(slide, l, t + h / 2 - Inches(0.2), w, Inches(0.4), "Diagram in images folder", 12, False, TEAL, PP_ALIGN.CENTER)
        return None
    return slide.shapes.add_picture(path, l, t, w, h)


def trap(slide, l, t, w, h, text):
    box(slide, l, t, w, h, TRAP_BG, ORANGE)
    tb(slide, l + Inches(0.12), t + Inches(0.06), w - Inches(0.2), Inches(0.28), "NEET TRAP", 11, True, ORANGE)
    tb(slide, l + Inches(0.12), t + Inches(0.32), w - Inches(0.2), h - Inches(0.38), text, 13, False, INK)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def s01(prs):
    s = new_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    rgb_fill(bg, CREAM)
    bg.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(2.35))
    rgb_fill(bar, TEAL_D)
    bar.line.fill.background()
    tb(s, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.9), "CELL: THE UNIT OF LIFE", 40, True, WHITE, PP_ALIGN.CENTER)
    tb(s, Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.45), "Class 11 Biology  ·  NCERT Chapter 8  ·  Morphology and organelles", 20, False, GOLD, PP_ALIGN.CENTER)
    tb(s, Inches(0.5), Inches(1.75), Inches(12.3), Inches(0.4), "ncert.nic.in/textbook/pdf/kebo108.pdf", 14, False, WHITE, PP_ALIGN.CENTER)
    picture(s, pick("gpt_cell_title.png", "cell_shapes.png"), Inches(0.35), Inches(2.5), Inches(8.3), Inches(4.45))
    box(s, Inches(8.8), Inches(2.5), Inches(4.15), Inches(4.45), TEAL)
    bullets(
        s,
        Inches(8.95),
        Inches(2.7),
        Inches(3.9),
        Inches(4.1),
        [
            "NCERT order 8.1 to 8.5.10",
            "Labelled size and shape",
            "Prokaryote vs eukaryote",
            "Plant vs animal cell",
            "NEET traps + PYQs",
            "21 slides · 16:9",
        ],
        16,
        WHITE,
        spacing=8,
    )
    footer(s, 1)
    notes(s, "NCERT Chapter 8 opener. Cell is the fundamental structural and functional unit of all living organisms. Unicellular vs multicellular. Open kebo108.pdf beside this deck.")
    return s


def s02(prs):
    s = new_slide(prs)
    header(s, "8.1  What is a cell?", "NCERT pp. opening of Ch. 8")
    bullets(
        s,
        Inches(0.4),
        Inches(1.05),
        Inches(7.1),
        Inches(5.7),
        [
            "Anything less than a complete cell does not ensure independent living.",
            "Cell = fundamental structural and functional unit of all living organisms.",
            "Unicellular: one cell does independent existence + essential functions of life.",
            "Multicellular: many cells (e.g. humans).",
            "Anton von Leeuwenhoek: first saw and described a live cell.",
            "Robert Brown: discovered the nucleus.",
            "Microscope → electron microscope revealed structural details.",
        ],
        18,
        spacing=10,
    )
    picture(s, pick("gpt_what_is_cell.png", "cell_shapes.png"), Inches(7.7), Inches(1.1), Inches(5.2), Inches(5.6))
    footer(s, 2)
    notes(s, "NCERT 8.1 WHAT IS A CELL? Unicellular organisms are capable of (i) independent existence and (ii) performing the essential functions of life.")
    return s


def s03(prs):
    s = new_slide(prs)
    header(s, "8.2  Cell theory", "Schleiden · Schwann · Virchow")
    picture(s, pick("gpt_cell_theory.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    timeline = [
        ("1838", "Matthias Schleiden", "German botanist. All plants are composed of different kinds of cells forming tissues."),
        ("1839", "Theodore Schwann", "British zoologist. Animal cells have a thin outer layer = plasma membrane. Cell wall unique to plant cells. Bodies of animals and plants are composed of cells and products of cells."),
        ("1855", "Rudolf Virchow", "Omnis cellula-e cellula: new cells from pre-existing cells. Gave cell theory its final shape."),
    ]
    x = 0.4
    for year, who, body in timeline:
        box(s, Inches(x), Inches(1.1), Inches(4.05), Inches(4.15), WHITE, TEAL)
        tb(s, Inches(x + 0.2), Inches(1.25), Inches(3.6), Inches(0.4), year, 22, True, ORANGE)
        tb(s, Inches(x + 0.2), Inches(1.7), Inches(3.6), Inches(0.7), who, 18, True, TEAL_D)
        tb(s, Inches(x + 0.2), Inches(2.45), Inches(3.65), Inches(2.6), body, 15, False, INK)
        x += 4.2
    box(s, Inches(0.4), Inches(5.45), Inches(12.5), Inches(1.5), SOFT, TEAL)
    tb(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(0.35), "Cell theory as understood today (NCERT)", 16, True, TEAL_D)
    tb(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.85), "(i) all living organisms are composed of cells and products of cells.\n(ii) all cells arise from pre-existing cells.", 18, True, INK)
    footer(s, 3)
    notes(s, "NCERT 8.2. Schleiden and Schwann formulated cell theory but did not explain how new cells formed. Virchow 1855 modified it.")
    return s


def s04(prs):
    s = new_slide(prs)
    header(s, "8.3  An overview of cell", "Onion peel vs cheek cell")
    bullets(
        s,
        Inches(0.4),
        Inches(1.0),
        Inches(7.3),
        Inches(4.4),
        [
            "Onion (plant): outer boundary = cell wall, then cell membrane.",
            "Human cheek (animal): outer delimiting structure = cell membrane.",
            "Nucleus: dense membrane-bound structure; chromosomes hold DNA.",
            "Eukaryotic = membrane-bound nucleus. Prokaryotic = no membrane-bound nucleus.",
            "Cytoplasm: semi-fluid matrix; main arena of cellular activities; keeps the cell in the living state.",
            "Eukaryotic membrane-bound organelles: ER, Golgi, lysosomes, mitochondria, microbodies, vacuoles.",
            "Prokaryotes lack those membrane-bound organelles.",
            "Ribosomes: non-membrane bound; in all cells; also in chloroplasts, mitochondria, on RER.",
            "Centrosome: another non-membrane organelle in animal cells; helps in cell division.",
        ],
        15,
        spacing=5,
    )
    picture(s, pick("gpt_cell_overview.png", "plant_animal.png"), Inches(7.85), Inches(1.0), Inches(5.1), Inches(4.5))
    trap(s, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.45), "Ribosomes occur in prokaryotes AND eukaryotes. Membrane-bound organelles listed above are eukaryotic only. Centrosome is animal (almost all plants lack centrioles).")
    footer(s, 4)
    notes(s, "NCERT 8.3 AN OVERVIEW OF CELL. Recollect onion peel and human cheek cells.")
    return s


def s05(prs):
    s = new_slide(prs)
    header(s, "Cell morphology: size", "NCERT 8.3 and Figure 8.2")
    picture(s, pick("gpt_cell_size.png", "cell_sizes.png"), Inches(0.35), Inches(0.95), Inches(8.3), Inches(5.95))
    box(s, Inches(8.8), Inches(0.95), Inches(4.15), Inches(5.95), WHITE, TEAL)
    tb(s, Inches(9.0), Inches(1.1), Inches(3.8), Inches(0.4), "NCERT numbers", 16, True, TEAL_D)
    bullets(
        s,
        Inches(8.95),
        Inches(1.55),
        Inches(3.85),
        Inches(5.1),
        [
            "Mycoplasma: 0.3 µm (smallest cells)",
            "Bacteria: 3 to 5 µm (typical also 1–2 µm in Fig. 8.2)",
            "PPLO: about 0.1 µm",
            "Viruses: 0.02–0.2 µm (shown for scale; not cells)",
            "Typical eukaryotic cell: 10–20 µm",
            "Human RBC: about 7.0 µm diameter",
            "Largest isolated single cell: ostrich egg",
            "Nerve cells: among the longest cells",
        ],
        14,
        spacing=6,
    )
    footer(s, 5)
    notes(s, "NCERT 8.3 and Figure 8.2. Confirm labels in kebo108.pdf Figure 8.2. Viruses are on the comparison figure; they are not cells.")
    return s


def s06(prs):
    s = new_slide(prs)
    header(s, "Cell morphology: shape", "NCERT 8.3 and Figure 8.1")
    picture(s, pick("gpt_cell_shapes.png", "cell_shapes.png"), Inches(0.3), Inches(0.92), Inches(8.6), Inches(5.55))
    box(s, Inches(9.05), Inches(0.95), Inches(3.9), Inches(5.5), WHITE, TEAL)
    tb(s, Inches(9.2), Inches(1.1), Inches(3.6), Inches(0.4), "NCERT shapes", 16, True, TEAL_D)
    bullets(
        s,
        Inches(9.15),
        Inches(1.55),
        Inches(3.7),
        Inches(4.7),
        [
            "Disc-like, polygonal, columnar, cuboid, thread-like, or irregular",
            "RBC: round and biconcave",
            "WBC: amoeboid",
            "Columnar epithelium: long and narrow",
            "Nerve cell: branched and long",
            "Mesophyll: round and oval",
            "Tracheid: elongated",
            "Shape may vary with the function they perform",
        ],
        14,
        spacing=6,
    )
    footer(s, 6)
    notes(s, "NCERT Figure 8.1. Confirm names on the official PDF: some OCR copies scramble WBC vs nerve-cell labels. Standard NCERT pairing: RBC biconcave; WBC amoeboid; columnar long and narrow; nerve branched and long; mesophyll round and oval; tracheid elongated.")
    return s


def s07(prs):
    s = new_slide(prs)
    header(s, "8.4  Prokaryotic cells", "Bacteria, cyanobacteria, mycoplasma / PPLO")
    picture(s, pick("gpt_bacterial_shapes.png", "bacterial_shapes.png"), Inches(0.3), Inches(0.95), Inches(7.6), Inches(3.55))
    bullets(
        s,
        Inches(8.05),
        Inches(0.95),
        Inches(4.9),
        Inches(3.6),
        [
            "Generally smaller; multiply more rapidly than eukaryotes",
            "No well-defined nucleus; genetic material is naked (nucleoid)",
            "Plasmids: small extra circular DNA (e.g. antibiotic resistance)",
            "No eukaryotic-type membrane-bound organelles",
            "70S ribosomes (50S + 30S)",
            "Inclusions: reserve material, not membrane-bound",
            "Mycoplasma: no cell wall",
        ],
        14,
        spacing=5,
    )
    trap(s, Inches(0.3), Inches(4.65), Inches(12.7), Inches(2.25), "Peroxisomes are eukaryotic. Nuclear envelope is absent. Ribosomes ARE present (70S). Mesosome is characteristic of prokaryotes (NCERT). Four shapes: bacillus (rod), coccus (spherical), vibrio (comma), spirillum (spiral).")
    footer(s, 7)
    notes(s, "NCERT 8.4. Represented by bacteria, blue-green algae, mycoplasma and PPLO. Organisation is fundamentally similar even though shapes vary.")
    return s


def s08(prs):
    s = new_slide(prs)
    header(s, "8.4.1  Cell envelope and surface", "Glycocalyx · wall · membrane · mesosome · flagella")
    picture(s, pick("gpt_prokaryote.png", "prokaryote_labelled.png"), Inches(0.25), Inches(0.92), Inches(7.7), Inches(6.05))
    for x, y, label_text in [
        (0.45, 1.35, "Glycocalyx"),
        (0.45, 2.05, "Cell wall"),
        (0.45, 2.75, "Plasma membrane"),
        (0.45, 3.45, "70S ribosomes"),
        (5.75, 1.75, "Nucleoid DNA"),
        (5.75, 2.55, "Plasmids"),
        (5.75, 3.35, "Mesosome"),
        (5.75, 4.15, "Flagellum"),
    ]:
        box(s, Inches(x), Inches(y), Inches(1.65), Inches(0.36), WHITE, TEAL)
        tb(s, Inches(x + 0.05), Inches(y + 0.04), Inches(1.55), Inches(0.25), label_text, 10, True, TEAL_D, PP_ALIGN.CENTER)
    bullets(
        s,
        Inches(8.1),
        Inches(0.95),
        Inches(4.9),
        Inches(4.4),
        [
            "Envelope: glycocalyx + cell wall + plasma membrane (one protective unit)",
            "Glycocalyx: slime (loose) or capsule (thick, tough)",
            "Wall determines shape; prevents bursting / collapsing",
            "Gram stain: Gram positive vs Gram negative",
            "Mesosome: membrane infoldings (vesicles, tubules, lamellae)",
            "Flagellum: filament + hook + basal body",
            "Pili: elongated tubular protein; fimbriae: bristle-like fibres",
            "Cyanobacteria: chromatophores (pigments) — not mesosomes",
        ],
        14,
        spacing=5,
    )
    trap(s, Inches(8.1), Inches(5.4), Inches(4.9), Inches(1.55), "Pili and fimbriae do NOT cause motility. Fimbriae attach to rocks / host tissues.")
    footer(s, 8)
    notes(s, "NCERT 8.4.1. Mesosome functions: cell wall formation, DNA replication and distribution, respiration, secretion, surface area. For NEET, use NCERT wording (not the later 'artifact' debate).")
    return s


def s09(prs):
    s = new_slide(prs)
    header(s, "8.5  Eukaryotic cells: plant vs animal", "NCERT Figure 8.3")
    picture(
        s,
        pick("gpt_plant_animal.png", "plant_animal.png"),
        Inches(0.25),
        Inches(0.92),
        Inches(8.5),
        Inches(6.05),
    )
    box(s, Inches(8.9), Inches(0.95), Inches(4.1), Inches(6.0), WHITE, TEAL)
    tb(s, Inches(9.05), Inches(1.1), Inches(3.8), Inches(0.35), "NCERT differences", 15, True, TEAL_D)
    bullets(
        s,
        Inches(9.0),
        Inches(1.5),
        Inches(3.85),
        Inches(5.2),
        [
            "Eukaryotes: protists, plants, animals, fungi",
            "Organised nucleus with envelope; chromosomes",
            "Plant only: cell wall, plastids, large central vacuole, plasmodesmata, middle lamella",
            "Animal typical: centrioles, lysosomes prominent, microvilli",
            "Centrioles absent in almost all plant cells",
            "Both: mitochondria, ER, Golgi, ribosomes, nucleus, peroxisomes",
        ],
        14,
        spacing=6,
    )
    footer(s, 9)
    notes(s, "NCERT 8.5 and Figure 8.3. Plant cell (a) and animal cell (b). Confirm organelle names on the textbook figure.")
    return s


def s10(prs):
    s = new_slide(prs)
    header(s, "8.5.1  Cell membrane", "Fluid mosaic · Singer and Nicolson 1972")
    picture(s, pick("gpt_membrane.png", "fluid_mosaic.png"), Inches(0.25), Inches(0.92), Inches(7.5), Inches(4.55))
    for x, y, label_text in [
        (0.55, 1.25, "Carbohydrate chain"),
        (1.0, 4.75, "Phospholipid bilayer"),
        (3.25, 1.25, "Integral protein"),
        (5.55, 4.75, "Cholesterol"),
    ]:
        box(s, Inches(x), Inches(y), Inches(1.75), Inches(0.38), WHITE, TEAL)
        tb(s, Inches(x + 0.05), Inches(y + 0.05), Inches(1.65), Inches(0.24), label_text, 10, True, TEAL_D, PP_ALIGN.CENTER)
    bullets(
        s,
        Inches(7.95),
        Inches(0.95),
        Inches(5.05),
        Inches(4.5),
        [
            "Lipids + proteins; major lipids = phospholipids in a bilayer",
            "Polar heads outside; hydrophobic tails inside",
            "Also cholesterol, protein, carbohydrate",
            "Human RBC membrane: ~52% protein, ~40% lipids",
            "Integral proteins: buried; peripheral: on surface",
            "Quasi-fluid lipids → lateral movement of proteins = fluidity",
            "Selectively permeable",
        ],
        14,
        spacing=5,
    )
    box(s, Inches(0.25), Inches(5.55), Inches(12.8), Inches(1.4), SOFT, TEAL)
    tb(
        s,
        Inches(0.4),
        Inches(5.65),
        Inches(12.5),
        Inches(1.2),
        "Passive transport: no energy; simple diffusion of neutral solutes; osmosis = water by diffusion. Polar molecules need carrier proteins. Active transport: against gradient, ATP, e.g. Na+/K+ pump.",
        16,
        False,
        INK,
    )
    footer(s, 10)
    notes(s, "NCERT 8.5.1 and Figure 8.4. Fluid mosaic model. Transport paragraph is a frequent long-answer (Exercise 6).")
    return s


def s11(prs):
    s = new_slide(prs)
    header(s, "8.5.2  Cell wall", "Fungi and plants")
    picture(s, pick("gpt_cell_wall.png", "cell_wall.png"), Inches(0.3), Inches(0.95), Inches(7.6), Inches(5.95))
    for x, y, label_text in [
        (0.5, 1.25, "Secondary wall"),
        (0.5, 2.0, "Primary wall"),
        (3.15, 3.15, "Middle lamella"),
        (5.65, 2.0, "Plasmodesma"),
        (5.65, 1.25, "Plasma membrane"),
    ]:
        box(s, Inches(x), Inches(y), Inches(1.8), Inches(0.4), WHITE, TEAL)
        tb(s, Inches(x + 0.05), Inches(y + 0.05), Inches(1.7), Inches(0.25), label_text, 10, True, TEAL_D, PP_ALIGN.CENTER)
    bullets(
        s,
        Inches(8.1),
        Inches(1.0),
        Inches(4.9),
        Inches(5.8),
        [
            "Non-living rigid covering over the plasma membrane",
            "Shape, mechanical protection, infection barrier, cell-to-cell interaction",
            "Algae: cellulose, galactans, mannans, calcium carbonate",
            "Other plants: cellulose, hemicellulose, pectins, proteins",
            "Primary wall: young cell, can grow",
            "Secondary wall: inner side (towards membrane) as cell matures",
            "Middle lamella: calcium pectate; glues neighbours",
            "Plasmodesmata: cytoplasmic connections through wall + middle lamella",
        ],
        15,
        spacing=7,
    )
    footer(s, 11)
    notes(s, "NCERT 8.5.2. Middle lamella calcium pectate is a frequent one-line NEET fact.")
    return s


def s12(prs):
    s = new_slide(prs)
    header(s, "8.5.3  Endomembrane system", "ER · Golgi · lysosomes · vacuoles")
    picture(s, pick("gpt_endomembrane.png", "endomembrane.png"), Inches(0.2), Inches(0.9), Inches(12.9), Inches(6.05))
    for x, y, label_text in [
        (0.55, 1.15, "Nucleus + RER"),
        (1.0, 4.65, "SER"),
        (5.25, 3.2, "Golgi: cis → trans"),
        (9.3, 2.05, "Lysosomes"),
        (9.15, 5.25, "Vacuole"),
    ]:
        box(s, Inches(x), Inches(y), Inches(2.05), Inches(0.42), WHITE, TEAL)
        tb(s, Inches(x + 0.05), Inches(y + 0.06), Inches(1.95), Inches(0.25), label_text, 11, True, TEAL_D, PP_ALIGN.CENTER)
    footer(s, 12)
    notes(s, "NCERT 8.5.3. Mitochondria, chloroplast and peroxisomes are NOT endomembrane. RER vs SER. Golgi cis/trans. Lysosomes = acid hydrolases from Golgi. Vacuole = tonoplast; up to 90% in plant cells.")
    return s


def s13(prs):
    s = new_slide(prs)
    header(s, "8.5.4–8.5.5  Mitochondria and plastids", "Double membrane · own DNA · 70S")
    picture(s, pick("gpt_energy_organelles.png", "energy_organelles.png"), Inches(0.2), Inches(0.9), Inches(12.9), Inches(6.05))
    for x, y, label_text in [
        (0.55, 1.15, "MITOCHONDRION"),
        (1.0, 5.75, "Cristae · matrix · circular DNA · 70S"),
        (7.25, 1.15, "CHLOROPLAST"),
        (7.45, 5.75, "Grana · stroma · circular DNA · 70S"),
    ]:
        box(s, Inches(x), Inches(y), Inches(5.25 if y > 5 else 2.4), Inches(0.48), WHITE, TEAL)
        tb(s, Inches(x + 0.05), Inches(y + 0.07), Inches(5.15 if y > 5 else 2.3), Inches(0.28), label_text, 12, True, TEAL_D, PP_ALIGN.CENTER)
    footer(s, 13)
    notes(s, "NCERT 8.5.4–8.5.5. Mitochondria: cristae, matrix, power houses, fission. Plastids: chloroplast, chromoplast (carotene, xanthophyll), leucoplast (amyloplast, elaioplast, aleuroplast). Chloroplast: grana, stroma, thylakoid lumen, stroma lamellae.")
    return s


def s14(prs):
    s = new_slide(prs)
    header(s, "8.5.10  Nucleus and chromosomes", "Robert Brown 1831 · Flemming chromatin")
    picture(s, pick("gpt_nucleus_chromosomes.png", "nucleus_chromosomes.png"), Inches(0.2), Inches(0.9), Inches(12.9), Inches(6.05))
    for x, y, label_text in [
        (0.55, 1.1, "Nuclear envelope + pores"),
        (1.25, 5.8, "Nucleolus · chromatin · nucleoplasm"),
        (7.35, 5.8, "Metacentric · sub-metacentric · acrocentric · telocentric"),
    ]:
        box(s, Inches(x), Inches(y), Inches(5.3), Inches(0.48), WHITE, TEAL)
        tb(s, Inches(x + 0.05), Inches(y + 0.07), Inches(5.2), Inches(0.28), label_text, 11, True, TEAL_D, PP_ALIGN.CENTER)
    footer(s, 14)
    notes(s, "NCERT 8.5.10. Envelope double membrane, perinuclear space 10–50 nm, pores for RNA and protein. Nucleolus: rRNA, not membrane-bound. Chromatin: DNA + histones + non-histones + RNA. Human cell ~2 m DNA in 46 chromosomes. Centromere types. Satellite = secondary constriction. Anucleate: mammalian RBC, sieve tubes.")
    return s


def s15(prs):
    s = new_slide(prs)
    header(s, "Recap tables", "Before the question block")
    picture(s, pick("gpt_recap_cells.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    # two comparison tables as text cards
    box(s, Inches(0.3), Inches(0.95), Inches(6.3), Inches(3.35), WHITE, TEAL)
    tb(s, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.35), "Prokaryote vs eukaryote", 16, True, TEAL_D)
    bullets(
        s,
        Inches(0.45),
        Inches(1.45),
        Inches(6.0),
        Inches(2.7),
        [
            "Nucleus: nucleoid vs envelope + chromosomes",
            "Organelles: no membrane-bound vs extensive",
            "Ribosomes: 70S vs 80S cytoplasm (70S in mito/plastid)",
            "Cell wall: present except mycoplasma; plants/fungi/algae",
            "Division: generally faster in prokaryotes",
        ],
        14,
        spacing=4,
    )
    box(s, Inches(6.8), Inches(0.95), Inches(6.2), Inches(3.35), WHITE, TEAL)
    tb(s, Inches(7.0), Inches(1.05), Inches(5.9), Inches(0.35), "Plant vs animal", 16, True, TEAL_D)
    bullets(
        s,
        Inches(6.95),
        Inches(1.45),
        Inches(5.9),
        Inches(2.7),
        [
            "Wall, plastids, large vacuole: plant",
            "Centrioles: animal (almost all plants lack them)",
            "Plasmodesmata / middle lamella: plant",
            "Both: nucleus, mitochondria, ER, Golgi, ribosomes",
        ],
        14,
        spacing=4,
    )
    box(s, Inches(0.3), Inches(4.45), Inches(12.7), Inches(2.45), SOFT, TEAL)
    tb(s, Inches(0.5), Inches(4.55), Inches(12.3), Inches(0.3), "One-line functions", 15, True, TEAL_D)
    tb(
        s,
        Inches(0.5),
        Inches(4.9),
        Inches(12.3),
        Inches(1.85),
        "Membrane: selective transport. Wall: shape + protection. RER: proteins. SER: lipids. Golgi: packaging / glycoproteins. Lysosome: acid hydrolases. Vacuole: sap / tonoplast. Mitochondrion: ATP. Chloroplast: photosynthesis. Ribosome: protein synthesis. Nucleus: control + heredity. Mesosome: wall, DNA, respiration in prokaryotes.",
        15,
        False,
        INK,
    )
    footer(s, 15)
    notes(s, "NCERT summary of Chapter 8. Use as a 3-minute oral recap before PYQs.")
    return s


def s16(prs):
    s = new_slide(prs)
    header(s, "NEET tricks (NCERT traps only)", "Slide 16 · keep the book line")
    picture(s, pick("gpt_neet_traps.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    traps = [
        "Endomembrane = ER + Golgi + lysosomes + vacuoles. Mitochondria, chloroplast, peroxisomes are OUT.",
        "Mesosome = plasma-membrane infolding (prokaryote). Chromatophores (pigments) = cyanobacteria. Do not mix.",
        "Flagella = motility. Pili and fimbriae ≠ motility. Fimbriae attach to rocks / host tissues.",
        "Inclusion bodies = reserve material, not membrane-bound. Polysome = ribosomes on one mRNA, not an inclusion.",
        "70S = 50S + 30S; 80S = 60S + 40S. Svedberg units do not add as ordinary numbers.",
        "Human RBC membrane ~52% protein and ~40% lipids.",
        "Mycoplasma: smallest cells (0.3 µm) and no cell wall.",
        "Centrioles absent in almost all plant cells. Wall, plastids, large vacuole absent in animal cells.",
        "Mature mammalian erythrocytes and sieve-tube cells lack a nucleus.",
        "Lysosomes: Golgi packaging; hydrolases active at acidic pH. Middle lamella = calcium pectate.",
        "Chromosomes: metacentric, sub-metacentric, acrocentric, telocentric. Satellite = secondary constriction.",
        "Ribosomes: all cells; also inside mitochondria and chloroplasts (70S) and on RER.",
    ]
    cols = [traps[:6], traps[6:]]
    for c, items in enumerate(cols):
        x = 0.3 + c * 6.5
        for i, text in enumerate(items):
            y = 0.95 + i * 0.95
            box(s, Inches(x), Inches(y), Inches(6.3), Inches(0.88), TRAP_BG, ORANGE)
            tb(s, Inches(x + 0.12), Inches(y + 0.08), Inches(6.05), Inches(0.72), text, 13, False, INK)
    footer(s, 16)
    notes(s, "All traps are restatements of NCERT 8.4–8.5.10. See QUESTIONS.md section D.")
    return s


def s17(prs):
    s = new_slide(prs)
    header(s, "Important questions — NCERT exercises", "Answers in notes and on slide 21")
    picture(s, pick("gpt_important_questions.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    box(s, Inches(0.3), Inches(0.9), Inches(8.8), Inches(6.05), WHITE, TEAL)
    items = [
        "Q2. New cells generate from?  (pre-existing cells)",
        "Q3. Match cristae / cisternae / thylakoids.",
        "Q5. What is a mesosome? Functions?",
        "Q6. Neutral solutes vs polar molecules across the membrane. Osmosis. Active transport.",
        "Q7. Two double-membrane organelles: structure, function, labelled diagrams.",
        "Q8. Characteristics of prokaryotic cells.",
        "Q11. Nuclear pores and their function.",
        "Q12. Lysosome vs vacuole (both endomembrane).",
        "Q13. Labelled diagrams: nucleus and centrosome.",
        "Q14. Centromere and chromosome types (draw).",
        "Exemplar: RER vs SER; plasmids; histones; group organelles by membrane number.",
        "Exercise 1 trap: Robert Brown did NOT discover the cell — he discovered the nucleus.",
    ]
    bullets(s, Inches(0.5), Inches(1.0), Inches(12.4), Inches(5.9), items, 16, spacing=6)
    footer(s, 17)
    notes(
        s,
        "Answers: Q2 (c) pre-existing cells. Q3 cristae=mitochondrial infoldings; cisternae=Golgi discs; thylakoids=flat sacs in stroma. Q5 see 8.4.1. Q6 simple diffusion vs carriers vs ATP. Q7 mitochondria + chloroplast. Q8 8.4. Q11 RNA and protein both ways. Q12 hydrolases vs tonoplast sap. Q14 meta/sub-meta/acro/telo. Source: ncert.nic.in/textbook/pdf/kebo108.pdf exercises; exemplar keep408.pdf.",
    )
    return s


def s18(prs):
    s = new_slide(prs)
    header(s, "Important MCQs (NCERT-style)", "Same stems NEET repeats")
    picture(s, pick("gpt_mcq.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    box(s, Inches(0.3), Inches(0.95), Inches(12.7), Inches(2.0), WHITE, TEAL)
    tb(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.3), "MCQ A  —  Match (NCERT Exercise 3 / frequent NEET matching)", 14, True, TEAL_D)
    tb(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.4), "Cristae · Cisternae · Thylakoids    vs    infoldings of mitochondria · disc-shaped sacs of Golgi · flat membranous sacs in stroma of chloroplast", 16, False, INK)

    box(s, Inches(0.3), Inches(3.1), Inches(6.2), Inches(3.75), WHITE, TEAL)
    tb(s, Inches(0.5), Inches(3.2), Inches(5.9), Inches(0.3), "MCQ B  —  Ribosomes", 14, True, TEAL_D)
    tb(s, Inches(0.5), Inches(3.55), Inches(5.9), Inches(3.1), "Eukaryotic cytoplasm: 80S (60S + 40S).\nProkaryote: 70S (50S + 30S).\nMitochondria and chloroplasts also have 70S.\nPalade (1953); RNA + protein; no membrane.\nS = Svedberg sedimentation coefficient.", 15, False, INK)

    box(s, Inches(6.7), Inches(3.1), Inches(6.3), Inches(3.75), WHITE, TEAL)
    tb(s, Inches(6.9), Inches(3.2), Inches(6.0), Inches(0.3), "MCQ C  —  Endomembrane", 14, True, TEAL_D)
    tb(s, Inches(6.9), Inches(3.55), Inches(6.0), Inches(3.1), "Includes: ER, Golgi, lysosomes, vacuoles.\nDoes NOT include: mitochondria, chloroplast, peroxisomes.\nSER is the major site for lipid synthesis, not carbohydrate as a primary NCERT line.", 15, False, INK)
    footer(s, 18)
    notes(s, "MCQ A answer: cristae–mitochondria; cisternae–Golgi; thylakoids–stroma sacs. MCQ B 80S=60+40; 70S=50+30. MCQ C ER+Golgi+lysosome+vacuole only. NCERT 8.5.3–8.5.6.")
    return s


def _mcq_block(s, l, t, w, h, title, body):
    box(s, Inches(l), Inches(t), Inches(w), Inches(h), WHITE, TEAL)
    tb(s, Inches(l + 0.15), Inches(t + 0.08), Inches(w - 0.3), Inches(0.32), title, 13, True, ORANGE)
    tb(s, Inches(l + 0.15), Inches(t + 0.4), Inches(w - 0.3), Inches(h - 0.5), body, 13, False, INK)


def s19(prs):
    s = new_slide(prs)
    header(s, "Previous-year questions — set A", "Year on each card · answers in notes")
    picture(s, pick("gpt_pyq_a.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    _mcq_block(
        s, 0.25, 0.95, 6.35, 2.9,
        "NEET 2023",
        "Statement I: In bacteria, mesosomes are formed by extensions of the plasma membrane.\nStatement II: Mesosomes help in DNA replication and cell wall formation.\n(1) I correct, II incorrect   (2) I incorrect, II correct\n(3) Both correct   (4) Both incorrect",
    )
    _mcq_block(
        s, 6.75, 0.95, 6.3, 2.9,
        "NEET 2020",
        "Which statement about inclusion bodies is incorrect?\n(1) Not bound by any membrane\n(2) Involved in ingestion of food particles\n(3) Lie free in the cytoplasm\n(4) Represent reserve material in cytoplasm",
    )
    _mcq_block(
        s, 0.25, 4.0, 6.35, 2.9,
        "NEET 2016",
        "Select the wrong statement.\n(1) Bacterial cell wall is made up of peptidoglycan\n(2) Pili and fimbriae are mainly involved in motility\n(3) Cyanobacteria lack flagellated cells\n(4) Mycoplasma is a wall-less microorganism",
    )
    _mcq_block(
        s, 6.75, 4.0, 6.3, 2.9,
        "NEET 2015",
        "Which structure is not found in a prokaryotic cell?\n(1) Mesosome   (2) Plasma membrane\n(3) Nuclear envelope   (4) Ribosome",
    )
    footer(s, 19)
    notes(
        s,
        "ANSWERS Set A: 2023 (3) both correct. 2020 (2) ingestion is wrong. 2016 (2) pili/fimbriae not motility. 2015 (3) nuclear envelope. Sources: smartachievers.online topic set; NCERT 8.4 / 8.4.1 / 8.4.2. Confirm on NTA paper PDFs.",
    )
    return s


def s20(prs):
    s = new_slide(prs)
    header(s, "Previous-year questions — set B", "Year on each card · answers in notes")
    picture(s, pick("gpt_pyq_b.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    _mcq_block(
        s, 0.25, 0.95, 6.35, 2.9,
        "NEET 2015",
        "Structures that help some bacteria attach to rocks and/or host tissues:\n(1) mesosomes   (2) holdfast\n(3) rhizoids   (4) fimbriae",
    )
    _mcq_block(
        s, 6.75, 0.95, 6.3, 2.9,
        "NEET 2015",
        "Which is not an inclusion body in prokaryotes?\n(1) Glycogen granule   (2) Polysome\n(3) Phosphate granule   (4) Cyanophycean granule",
    )
    _mcq_block(
        s, 0.25, 4.0, 6.35, 2.9,
        "NEET 2014  /  NEET 2024",
        "2014: Which structures perform the function of mitochondria in bacteria?\n(1) Nucleoid (2) Ribosomes (3) Cell wall (4) Mesosomes\n\n2024: Mesosome in a cell is a special structure formed by the extension of plasma membrane (not a polysome).",
    )
    _mcq_block(
        s, 6.75, 4.0, 6.3, 2.9,
        "NEET 2025",
        "A, B, C true: eukaryotic ribosomes 80S, prokaryotic 70S; each has two subunits; 80S = 60S+40S and 70S = 50S+30S.\nAlso 2025: specialised membranous structure for wall formation, DNA replication and respiration = mesosome (not chromatophore).",
    )
    footer(s, 20)
    notes(
        s,
        "ANSWERS Set B: 2015 fimbriae (4). 2015 polysome not inclusion (2). 2014 mesosomes (4). 2024 mesosome = membrane extension. 2025 A,B,C true; specialised structure = mesosome. Sources: smartachievers.online; EduRev 2016–2025 list for 2024 wording. NCERT 8.4.1, 8.4.2, 8.5.6.",
    )
    return s


def s21(prs):
    s = new_slide(prs)
    header(s, "Answer key and where to check", "Do not skip the textbook figure")
    picture(s, pick("gpt_answer_key.png"), Inches(0.15), Inches(0.82), Inches(13.03), Inches(6.3))
    box(s, Inches(0.3), Inches(0.9), Inches(8.7), Inches(6.05), WHITE, TEAL)
    bullets(
        s,
        Inches(0.4),
        Inches(1.0),
        Inches(12.5),
        Inches(5.9),
        [
            "Slide 17: Ex.1 (a) wrong — Brown found the nucleus. Ex.2 (c). Ex.3 cristae–mito, cisternae–Golgi, thylakoids–stroma. Ex.4 (c) prokaryotes lack membrane-bound organelles.",
            "Slide 18: Matching as above. 80S=60S+40S; 70S=50S+30S. Endomembrane excludes mito, chloroplast, peroxisome.",
            "Slide 19: 2023 both statements correct. 2020 inclusion ≠ food ingestion. 2016 pili/fimbriae ≠ motility. 2015 no nuclear envelope in prokaryote.",
            "Slide 20: 2015 fimbriae attach. 2015 polysome ≠ inclusion. 2014 mesosome ~ mitochondrial function in bacteria. 2025 80S/70S subunits A,B,C.",
            "Full options + URLs: QUESTIONS.md in this folder.",
            "Official chapter: https://ncert.nic.in/textbook/pdf/kebo108.pdf",
            "Exemplar: https://ncert.nic.in/pdf/publication/exemplarproblem/classXI/biology/keep408.pdf",
            "Papers: nta.ac.in  ·  Image sources: IMAGE-AND-NCERT-REFERENCES.md",
        ],
        15,
        spacing=7,
    )
    footer(s, 21)
    notes(s, "Close class by opening kebo108.pdf Figures 8.1–8.13. All PYQ years should be confirmed on the NTA PDF before a printed test.")
    return s


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    for fn in (s01, s02, s03, s04, s05, s06, s07, s08, s09, s10, s11, s12, s13, s14, s15, s16, s17, s18, s19, s20, s21):
        fn(prs)
    prs.save(str(OUT))
    print("Wrote", OUT, "slides:", len(prs.slides))


if __name__ == "__main__":
    main()
